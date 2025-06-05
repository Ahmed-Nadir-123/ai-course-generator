"""
AI Course Generator Backend
Flask server with Gemini AI integration for comprehensive course generation
"""

import os
import time
import logging
from datetime import datetime
from pathlib import Path
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from dotenv import load_dotenv

# Load environment variables - prioritize .env.local for development
if os.path.exists('.env.local'):
    load_dotenv('.env.local')
else:
    load_dotenv('.env')

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Create Flask app
app = Flask(__name__)
CORS(app)

# Global variables for rate limiting
last_request_time = 0
request_count_per_minute = 0
minute_start_time = time.time()

# Create directories
UPLOAD_FOLDER = 'generated_files'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

print("Starting AI Course Generator backend...")
print(f"Files directory: {UPLOAD_FOLDER}")
print("Server will be available at http://localhost:5000")

def rate_limit_decorator(func):
    """Decorator to implement rate limiting for API calls"""
    def wrapper(*args, **kwargs):
        global last_request_time, request_count_per_minute, minute_start_time
        
        current_time = time.time()
        
        # Reset counter every minute
        if current_time - minute_start_time > 60:
            request_count_per_minute = 0
            minute_start_time = time.time()
        
        # Minimum delay between requests (4 seconds for free tier)
        time_since_last = current_time - last_request_time
        if time_since_last < 4:
            time.sleep(4 - time_since_last)
        
        request_count_per_minute += 1
        last_request_time = time.time()
        
        return func(*args, **kwargs)
    return wrapper

@rate_limit_decorator
def get_gemini_response(user_message: str) -> dict:
    """
    Calls Gemini AI with a structured prompt and returns the parsed response.
    Implements rate limiting and fallback model strategy.
    """
    SYSTEM_PROMPT = '''You are an expert educational designer and course creator with deep expertise in curriculum development, instructional design, and modern teaching methodologies. 

Create a comprehensive, professional-grade course that follows industry best practices and educational standards. Generate detailed course content based on user requests with the following structure:

1. COURSE OVERVIEW
   - Title (engaging and descriptive)
   - Description (compelling 2-3 sentence overview)
   - Learning Objectives (specific, measurable, achievable)
   - Prerequisites (clear requirements)
   - Estimated Duration (realistic timeframe)
   - Target Audience (who should take this course)

2. COURSE DELIVERY PLAN
   - Timeline (week-by-week breakdown)
   - Learning Path (logical progression)
   - Assessment Methods (varied and comprehensive)
   - Delivery Format (online, hybrid, in-person considerations)

3. DETAILED MODULES
   For each module (aim for 6-8 modules):
   - Module Title (clear and engaging)
   - Description (what students will learn)
   - Key Topics (comprehensive list)
   - Learning Outcomes (specific skills/knowledge gained)
   - PowerPoint Content (detailed slide outlines with bullet points)
   - Lab Activities/Projects (hands-on practical work)
   - Resources (books, articles, tools, websites)
   - Assessment (quizzes, assignments, projects)

4. PRACTICAL COMPONENTS
   - Hands-on Exercises (step-by-step instructions)
   - Real-world Projects (portfolio-worthy work)
   - Assessment Criteria (detailed rubrics)
   - Capstone Project (comprehensive final project)

5. SUPPLEMENTARY MATERIALS
   - Additional Resources
   - Further Reading
   - Industry Tools and Software
   - Community and Support

Make the content:
- Engaging and interactive
- Industry-relevant and current
- Practical and applicable
- Well-structured and progressive
- Suitable for both beginners and intermediate learners

Format the response in a clear, structured manner suitable for:
- Display in chat
- PowerPoint generation
- PDF creation
- Lab instruction sheets'''
    
    # List of models to try (from most preferred to fallback)
    # Using the most powerful and latest models for best course generation
    models_to_try = [
        "models/gemini-2.5-pro-preview-05-06",  # Latest and most powerful Gemini 2.5 Pro
        "models/gemini-2.0-flash-exp",  # Latest experimental Gemini 2.0
        "models/gemini-2.0-flash",  # Stable Gemini 2.0
        "models/gemini-1.5-pro-latest",  # Latest Gemini 1.5 Pro
        "models/gemini-1.5-flash-latest",  # Fallback to efficient model
        "models/gemini-1.5-flash"  # Final fallback
    ]
    
    try:
        # Configure API key from environment variables
        api_key = os.getenv("GOOGLE_API_KEY") or os.getenv("GEMINI_API_KEY")
        if not api_key:
            return {"success": False, "error": "Gemini API key not found. Please set GOOGLE_API_KEY or GEMINI_API_KEY environment variable."}
        
        genai.configure(api_key=api_key)
        
        # Try each model in order of preference
        for model_name in models_to_try:
            try:
                logger.info(f"Trying model: {model_name}")
                model = genai.GenerativeModel(model_name)
                
                # Create the full prompt
                full_prompt = f"{SYSTEM_PROMPT}\n\nUser Request: {user_message}"
                
                # Generate response with retry logic
                response = model.generate_content(
                    full_prompt,
                    generation_config=genai.types.GenerationConfig(
                        temperature=0.7,
                        max_output_tokens=8192,
                        top_p=0.9,
                        top_k=40
                    )
                )
                
                if response and response.text:
                    logger.info(f"Successfully generated response using {model_name}")
                    return {
                        "success": True, 
                        "response": response.text,
                        "model_used": model_name
                    }
                    
            except Exception as model_error:
                logger.warning(f"Model {model_name} failed: {str(model_error)}")
                
                # Handle specific quota errors
                if "quota" in str(model_error).lower() or "rate limit" in str(model_error).lower():
                    continue  # Try next model
                elif "400" in str(model_error):
                    continue  # Bad request, try next model
                else:
                    # For other errors, try the next model
                    continue
        
        # If all models failed, return a helpful error
        return {
            "success": False, 
            "error": "All AI models are currently unavailable. This might be due to rate limits or quota restrictions.",
            "suggestion": "Please try again in a few minutes, or check your API quota."
        }
        
    except Exception as e:
        logger.error(f"Gemini AI error: {e}")
        return {"success": False, "error": f"AI service error: {str(e)}"}

# Flask routes

@app.route('/')
def home():
    return jsonify({"message": "AI Course Generator API is running", "status": "active"})

@app.route('/chat', methods=['POST'])
def chat():
    try:
        data = request.get_json()
        if not data or 'message' not in data:
            return jsonify({"success": False, "error": "Message is required"}), 400
        
        user_message = data['message']
        session_id = data.get('session_id', 'default')
        
        # Get response from Gemini AI
        result = get_gemini_response(user_message)
        
        # Log the interaction
        logger.info(f"Session {session_id}: User asked about '{user_message[:50]}...'")
        
        if result["success"]:
            return jsonify(result)
        else:
            return jsonify(result), 500
            
    except Exception as e:
        logger.error(f"Chat endpoint error: {e}")
        return jsonify({"success": False, "error": str(e)}), 500

def create_powerpoint(course_content: str, filename: str) -> bool:
    """Create a PowerPoint presentation from course content"""
    try:
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "AI Generated Course"
        subtitle.text = "Comprehensive Learning Program"
        
        # Content slides - split content by lines and create slides
        lines = course_content.split('\n')
        current_slide = None
        current_content = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # Check if this is a major heading (Module, Section, etc.)
            if (line.startswith('#') or 
                line.upper().startswith('MODULE') or 
                line.upper().startswith('SECTION') or
                line.upper().startswith('CHAPTER')):
                
                # Save previous slide if exists
                if current_slide and current_content:
                    content_text = '\n'.join(current_content)
                    if len(current_slide.shapes) > 1:
                        current_slide.shapes[1].text = content_text
                
                # Create new slide
                bullet_slide_layout = prs.slide_layouts[1]
                current_slide = prs.slides.add_slide(bullet_slide_layout)
                title = current_slide.shapes.title
                title.text = line.replace('#', '').strip()
                current_content = []
                
            else:
                # Add to current slide content
                current_content.append(line)
        
        # Don't forget the last slide
        if current_slide and current_content:
            content_text = '\n'.join(current_content)
            if len(current_slide.shapes) > 1:
                current_slide.shapes[1].text = content_text
        
        # Save presentation
        filepath = os.path.join(UPLOAD_FOLDER, filename)
        prs.save(filepath)
        return True
        
    except Exception as e:
        logger.error(f"PowerPoint creation error: {e}")
        return False

def create_pdf(course_content: str, filename: str) -> bool:
    """Create a PDF document from course content"""
    try:
        filepath = os.path.join(UPLOAD_FOLDER, filename)
        doc = SimpleDocTemplate(filepath, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            spaceAfter=30,
            textColor=colors.darkblue
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=16,
            spaceAfter=12,
            textColor=colors.darkgreen
        )
        
        # Add title
        story.append(Paragraph("AI Generated Course Content", title_style))
        story.append(Spacer(1, 20))
        
        # Process content
        lines = course_content.split('\n')
        for line in lines:
            line = line.strip()
            if not line:
                story.append(Spacer(1, 6))
                continue
            
            # Format based on content type
            if (line.startswith('#') or 
                line.upper().startswith('MODULE') or 
                line.upper().startswith('SECTION')):
                story.append(Paragraph(line.replace('#', '').strip(), heading_style))
            else:
                story.append(Paragraph(line, styles['Normal']))
            
            story.append(Spacer(1, 6))
        
        # Build PDF
        doc.build(story)
        return True
        
    except Exception as e:
        logger.error(f"PDF creation error: {e}")
        return False

@app.route('/generate_ppt', methods=['POST'])
def generate_ppt():
    try:
        data = request.get_json()
        if not data or 'course_content' not in data:
            return jsonify({"success": False, "error": "Course content is required"}), 400
        
        course_content = data['course_content']
        timestamp = datetime.now().strftime('%Y%m%d%H%M%S')
        filename = f"course_{timestamp}.pptx"
        
        if create_powerpoint(course_content, filename):
            return jsonify({
                "success": True,
                "filename": filename,
                "download_url": f"/download/{filename}"
            })
        else:
            return jsonify({"success": False, "error": "Failed to create PowerPoint"}), 500
            
    except Exception as e:
        logger.error(f"PPT generation error: {e}")
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/generate_pdf', methods=['POST'])
def generate_pdf():
    try:
        data = request.get_json()
        if not data or 'course_content' not in data:
            return jsonify({"success": False, "error": "Course content is required"}), 400
        
        course_content = data['course_content']
        timestamp = datetime.now().strftime('%Y%m%d%H%M%S')
        filename = f"course_{timestamp}.pdf"
        
        if create_pdf(course_content, filename):
            return jsonify({
                "success": True,
                "filename": filename,
                "download_url": f"/download/{filename}"
            })
        else:
            return jsonify({"success": False, "error": "Failed to create PDF"}), 500
            
    except Exception as e:
        logger.error(f"PDF generation error: {e}")
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/download/<filename>')
def download_file(filename):
    try:
        filepath = os.path.join(UPLOAD_FOLDER, filename)
        if os.path.exists(filepath):
            return send_file(filepath, as_attachment=True)
        else:
            return jsonify({"error": "File not found"}), 404
    except Exception as e:
        logger.error(f"Download error: {e}")
        return jsonify({"error": str(e)}), 500

@app.route('/files')
def list_files():
    try:
        files = []
        for filename in os.listdir(UPLOAD_FOLDER):
            filepath = os.path.join(UPLOAD_FOLDER, filename)
            if os.path.isfile(filepath):
                files.append({
                    "name": filename,
                    "size": os.path.getsize(filepath),
                    "modified": os.path.getmtime(filepath)
                })
        return jsonify({"files": files})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)
